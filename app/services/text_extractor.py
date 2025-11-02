"""
Service for extracting text from structured document formats without OCR.

For file types that contain structured text (HTML, CSV, DOCX, XLSX, etc.),
this service extracts the text directly and converts to JSON format for schema extraction.
"""

import csv
import json
import logging
from io import BytesIO, StringIO
from typing import Any

logger = logging.getLogger(__name__)


class TextExtractionService:
    """Extract text from structured documents without OCR"""
    
    @staticmethod
    def can_extract_text(content_type: str | None, filename: str) -> bool:
        """
        Determine if we can extract text directly from this file type
        
        Returns True for HTML, CSV, DOCX, XLSX, PPTX, TXT files
        Returns False for images and PDFs (which need OCR)
        """
        if not content_type and not filename:
            return False
        
        filename_lower = filename.lower() if filename else ""
        
        # File types that contain extractable text
        text_formats = [
            '.csv', '.html', '.htm', '.txt',
            '.docx', '.xlsx', '.pptx',
            '.json', '.xml'
        ]
        
        # Check by extension
        for ext in text_formats:
            if filename_lower.endswith(ext):
                return True
        
        # Check by MIME type
        if content_type:
            text_mime_types = [
                'text/csv',
                'text/html',
                'text/plain',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'application/json',
                'text/xml',
                'application/xml'
            ]
            if any(mime in content_type for mime in text_mime_types):
                return True
        
        return False
    
    @staticmethod
    def extract_text(file_bytes: bytes, content_type: str | None, filename: str) -> dict[str, Any]:
        """
        Extract text and structure from document
        
        Returns:
            dict with:
                - text: extracted text content
                - structured_data: parsed structured data (for CSV, JSON, etc.)
                - format: detected format type
        """
        filename_lower = filename.lower() if filename else ""
        
        try:
            # CSV files
            if filename_lower.endswith('.csv') or (content_type and 'text/csv' in content_type):
                return TextExtractionService._extract_csv(file_bytes, filename)
            
            # HTML files
            elif filename_lower.endswith(('.html', '.htm')) or (content_type and 'text/html' in content_type):
                return TextExtractionService._extract_html(file_bytes, filename)
            
            # Plain text
            elif filename_lower.endswith('.txt') or (content_type and 'text/plain' in content_type):
                return TextExtractionService._extract_text(file_bytes, filename)
            
            # JSON files
            elif filename_lower.endswith('.json') or (content_type and 'application/json' in content_type):
                return TextExtractionService._extract_json(file_bytes, filename)
            
            # DOCX files
            elif filename_lower.endswith('.docx'):
                return TextExtractionService._extract_docx(file_bytes, filename)
            
            # XLSX files
            elif filename_lower.endswith('.xlsx'):
                return TextExtractionService._extract_xlsx(file_bytes, filename)
            
            # PPTX files
            elif filename_lower.endswith('.pptx'):
                return TextExtractionService._extract_pptx(file_bytes, filename)
            
            else:
                logger.warning(f"Unsupported format for text extraction: {filename}")
                return {
                    "text": "",
                    "structured_data": None,
                    "format": "unknown",
                    "error": "Unsupported format"
                }
        
        except Exception as e:
            logger.error(f"Text extraction failed for {filename}: {str(e)}")
            return {
                "text": "",
                "structured_data": None,
                "format": "error",
                "error": str(e)
            }
    
    @staticmethod
    def _extract_csv(file_bytes: bytes, filename: str) -> dict[str, Any]:
        """Extract data from CSV file"""
        try:
            csv_text = file_bytes.decode('utf-8')
            csv_reader = csv.DictReader(StringIO(csv_text))
            rows = list(csv_reader)
            
            # Convert to JSON structure
            return {
                "text": csv_text,
                "structured_data": {
                    "rows": rows,
                    "row_count": len(rows),
                    "columns": list(rows[0].keys()) if rows else []
                },
                "format": "csv"
            }
        except Exception as e:
            logger.error(f"CSV extraction failed: {e}")
            # Fallback: just return raw text
            try:
                return {
                    "text": file_bytes.decode('utf-8'),
                    "structured_data": None,
                    "format": "csv",
                    "error": str(e)
                }
            except:
                return {"text": "", "structured_data": None, "format": "csv", "error": str(e)}
    
    @staticmethod
    def _extract_html(file_bytes: bytes, filename: str) -> dict[str, Any]:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            html_content = file_bytes.decode('utf-8')
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return {
                "text": text,
                "structured_data": {
                    "title": soup.title.string if soup.title else None,
                    "html": html_content
                },
                "format": "html"
            }
        except ImportError:
            # Fallback without BeautifulSoup
            logger.warning("BeautifulSoup not available, using raw HTML")
            try:
                return {
                    "text": file_bytes.decode('utf-8'),
                    "structured_data": None,
                    "format": "html"
                }
            except:
                return {"text": "", "structured_data": None, "format": "html"}
        except Exception as e:
            logger.error(f"HTML extraction failed: {e}")
            return {"text": "", "structured_data": None, "format": "html", "error": str(e)}
    
    @staticmethod
    def _extract_text(file_bytes: bytes, filename: str) -> dict[str, Any]:
        """Extract text from plain text file"""
        try:
            text = file_bytes.decode('utf-8')
            return {
                "text": text,
                "structured_data": None,
                "format": "text"
            }
        except:
            # Try other encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    text = file_bytes.decode(encoding)
                    return {
                        "text": text,
                        "structured_data": None,
                        "format": "text"
                    }
                except:
                    continue
            return {"text": "", "structured_data": None, "format": "text", "error": "Encoding error"}
    
    @staticmethod
    def _extract_json(file_bytes: bytes, filename: str) -> dict[str, Any]:
        """Extract data from JSON file"""
        try:
            json_text = file_bytes.decode('utf-8')
            json_data = json.loads(json_text)
            
            return {
                "text": json.dumps(json_data, indent=2),
                "structured_data": json_data,
                "format": "json"
            }
        except Exception as e:
            logger.error(f"JSON extraction failed: {e}")
            return {"text": "", "structured_data": None, "format": "json", "error": str(e)}
    
    @staticmethod
    def _extract_docx(file_bytes: bytes, filename: str) -> dict[str, Any]:
        """Extract text from DOCX file"""
        try:
            from docx import Document
            doc = Document(BytesIO(file_bytes))
            
            # Extract all paragraphs
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = '\n\n'.join(paragraphs)
            
            # Extract tables
            tables_data = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables_data.append(table_data)
            
            return {
                "text": text,
                "structured_data": {
                    "paragraphs": paragraphs,
                    "tables": tables_data if tables_data else None
                },
                "format": "docx"
            }
        except ImportError:
            logger.warning("python-docx not installed, cannot extract DOCX")
            return {"text": "", "structured_data": None, "format": "docx", "error": "python-docx not installed"}
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return {"text": "", "structured_data": None, "format": "docx", "error": str(e)}
    
    @staticmethod
    def _extract_xlsx(file_bytes: bytes, filename: str) -> dict[str, Any]:
        """Extract data from XLSX file"""
        try:
            import openpyxl
            from openpyxl import load_workbook
            
            wb = load_workbook(BytesIO(file_bytes), data_only=True)
            
            sheets_data = {}
            all_text = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    if any(row_data):  # Skip empty rows
                        sheet_data.append(row_data)
                        all_text.append(', '.join(row_data))
                
                sheets_data[sheet_name] = sheet_data
            
            return {
                "text": '\n'.join(all_text),
                "structured_data": {
                    "sheets": sheets_data,
                    "sheet_names": wb.sheetnames
                },
                "format": "xlsx"
            }
        except ImportError:
            logger.warning("openpyxl not installed, cannot extract XLSX")
            return {"text": "", "structured_data": None, "format": "xlsx", "error": "openpyxl not installed"}
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            return {"text": "", "structured_data": None, "format": "xlsx", "error": str(e)}
    
    @staticmethod
    def _extract_pptx(file_bytes: bytes, filename: str) -> dict[str, Any]:
        """Extract text from PPTX file"""
        try:
            from pptx import Presentation
            
            prs = Presentation(BytesIO(file_bytes))
            
            slides_text = []
            all_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)
                        all_text.append(shape.text)
                
                slides_text.append({
                    "slide_number": slide_num,
                    "content": slide_content
                })
            
            return {
                "text": '\n\n'.join(all_text),
                "structured_data": {
                    "slides": slides_text,
                    "slide_count": len(slides_text)
                },
                "format": "pptx"
            }
        except ImportError:
            logger.warning("python-pptx not installed, cannot extract PPTX")
            return {"text": "", "structured_data": None, "format": "pptx", "error": "python-pptx not installed"}
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return {"text": "", "structured_data": None, "format": "pptx", "error": str(e)}
